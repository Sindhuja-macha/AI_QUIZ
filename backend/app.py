import os
import json
import fitz  # PyMuPDF
from flask import Flask, request, jsonify
from flask_cors import CORS
from pptx import Presentation
from dotenv import load_dotenv
from google import genai
from google.genai import types
from pydantic import BaseModel, Field
from typing import List, Dict

# --------------------------------------------------
# Load Environment Variables
# --------------------------------------------------
base_dir = os.path.dirname(os.path.abspath(__file__))
load_dotenv(os.path.join(base_dir, ".env"))

app = Flask(__name__)

# Configure CORS to accept requests from any deployed frontend link
CORS(app, resources={r"/api/*": {"origins": "*"}})

# --------------------------------------------------
# Gemini Initialization
# --------------------------------------------------
gemini_key = os.getenv("GEMINI_API_KEY")

if not gemini_key or "your_" in gemini_key.lower():
    print("\n⚠️ GEMINI_API_KEY is missing in backend/.env")
else:
    print(f"✅ Gemini API Loaded Successfully.")

client = genai.Client(api_key=gemini_key)

# --------------------------------------------------
# Pydantic Schema for Reliable JSON Output
# --------------------------------------------------
class QuestionModel(BaseModel):
    id: int
    topic: str = Field(description="The core subject area or category of this question (e.g., 'Prompt Engineering', 'Database Management', 'AgTech')")
    question: str
    options: Dict[str, str] = Field(description="Dictionary containing exactly 4 keys: A, B, C, D")
    correctAnswer: str = Field(description="Must match exactly one string key from options: 'A', 'B', 'C', or 'D'")
    explanation: str

class QuizModel(BaseModel):
    questions: List[QuestionModel]

# --------------------------------------------------
# PPT Parser API
# --------------------------------------------------
@app.route("/api/parse-ppt", methods=["POST"])
def parse_ppt():
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded."}), 400

    file = request.files["file"]
    if not file.filename.endswith(".pptx"):
        return jsonify({"error": "Only .pptx files are allowed."}), 400

    try:
        prs = Presentation(file)
        extracted_text = []

        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)

            if slide_text:
                extracted_text.append(f"Slide {i+1}:\n" + "\n".join(slide_text))

        full_text = "\n\n".join(extracted_text)

        return jsonify({
            "slideCount": len(prs.slides),
            "wordCount": len(full_text.split()),
            "previewText": full_text[:1000] + ("..." if len(full_text) > 1000 else ""),
            "fullText": full_text,
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# --------------------------------------------------
# PDF Parser API
# --------------------------------------------------
@app.route("/api/parse-pdf", methods=["POST"])
def parse_pdf():
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded."}), 400

    file = request.files["file"]
    if not file.filename.endswith(".pdf"):
        return jsonify({"error": "Only PDF files are allowed."}), 400

    try:
        pdf = fitz.open(stream=file.read(), filetype="pdf")
        extracted_text = ""

        for page in pdf:
            extracted_text += page.get_text()

        return jsonify({
            "pageCount": len(pdf),
            "wordCount": len(extracted_text.split()),
            "previewText": extracted_text[:1000] + ("..." if len(extracted_text) > 1000 else ""),
            "fullText": extracted_text,
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# --------------------------------------------------
# Quiz Generator API
# --------------------------------------------------
@app.route("/api/generate-quiz", methods=["POST"])
def generate_quiz():
    data = request.json or {}
    text_content = data.get("text")
    count = data.get("count", 10)
    difficulty = data.get("difficulty", "Medium")

    if not text_content:
        return jsonify({"error": "No study material found."}), 400

    prompt = f"""
You are an expert educator. Your task is to analyze the provided study material and generate a structured quiz assessment.

Generate exactly {count} multiple choice questions.
Target Difficulty Level: {difficulty}

Core Requirements:
1. Thoroughly parse the text and extract precise technical classifications or logical sub-topics for the 'topic' field.
2. Ensure realistic distractor options that test core semantic understanding.
3. Use conceptual and scenario-based questions over surface-level verbatim lookups.
4. Provide highly informative summaries in the 'explanation' property that act as mini learning items.

Study Material Context:
{text_content}
"""

    try:
        # Enforce structural integrity via Google GenAI type settings
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt,
            config=types.GenerateContentConfig(
                response_mime_type="application/json",
                response_schema=QuizModel,
                temperature=0.3
            ),
        )

        # The SDK guarantees that response.text strictly maps to our QuizModel definition
        parsed_response = json.loads(response.text)
        return jsonify(parsed_response.get("questions", []))

    except Exception as e:
        return jsonify({"error": f"Failed to generate structured quiz content: {str(e)}"}), 500

# --------------------------------------------------
# Production/Development Execution Flow
# --------------------------------------------------
if __name__ == "__main__":
    # Determine execution context safely
    is_debug = os.getenv("FLASK_ENV") == "development" or os.getenv("DEBUG", "False").lower() in ("true", "1")
    port = int(os.getenv("PORT", 5001))
    
    print(f"\n🚀 Starting Flask Backend Engine...")
    print(f"📍 Mode: {'Development' if is_debug else 'Production'}")
    print(f"🔗 Target: http://localhost:{port}")
    
    # use_reloader=False stops the infinite file system checking loops on Windows instances
    app.run(host="0.0.0.0", port=port, debug=is_debug, use_reloader=False)